"""Generate the AgriChain judge presentation → presentation/AgriChain.pptx.

Run:  python -m presentation.make_deck
"""
from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

from config import ARCHITECTURE_PNG, FEATURE_IMPORTANCE_PNG, PRESENTATION_DIR  # noqa: E402

GREEN = RGBColor(0x2E, 0x7D, 0x32)
DARK = RGBColor(0x1B, 0x1B, 0x1B)
GREY = RGBColor(0x45, 0x5A, 0x64)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE8, 0xF5, 0xE9)

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
DECK_PATH = PRESENTATION_DIR / "AgriChain.pptx"


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _band(slide, color=GREEN, height=Inches(1.15)):
    box = slide.shapes.add_shape(1, 0, 0, EMU_W, height)  # 1 = rectangle
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def _text(slide, left, top, width, height, text, size, color=DARK,
          bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def title_slide(prs, title, subtitle, footer):
    s = _blank(prs)
    box = s.shapes.add_shape(1, 0, Inches(2.4), EMU_W, Inches(2.7))
    box.fill.solid(); box.fill.fore_color.rgb = GREEN
    box.line.fill.background(); box.shadow.inherit = False
    _text(s, Inches(0.8), Inches(2.75), Inches(11.7), Inches(1.2),
          title, 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.9),
          subtitle, 22, WHITE, align=PP_ALIGN.CENTER)
    _text(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6),
          footer, 14, GREY, align=PP_ALIGN.CENTER)
    return s


def bullets_slide(prs, title, bullets):
    s = _blank(prs)
    _band(s)
    _text(s, Inches(0.6), Inches(0.28), Inches(12.1), Inches(0.7),
          title, 30, WHITE, bold=True)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = ("• " if level == 0 else "– ") + text
        run.font.size = Pt(22 if level == 0 else 18)
        run.font.color.rgb = DARK if level == 0 else GREY
        run.font.bold = level == 0
        p.space_after = Pt(8)
    return s


def image_slide(prs, title, image_path, caption=""):
    s = _blank(prs)
    _band(s)
    _text(s, Inches(0.6), Inches(0.28), Inches(12.1), Inches(0.7),
          title, 30, WHITE, bold=True)
    if Path(image_path).exists():
        s.shapes.add_picture(str(image_path), Inches(1.6), Inches(1.45),
                             height=Inches(5.0))
    else:
        _text(s, Inches(1.0), Inches(3.0), Inches(11), Inches(1),
              f"[missing image: {image_path}]", 18, GREY)
    if caption:
        _text(s, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.6),
              caption, 14, GREY, align=PP_ALIGN.CENTER)
    return s


def build() -> Path:
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H

    title_slide(
        prs,
        "AgriChain",
        "Blockchain-Based Agricultural Supply Chain Transparency",
        "Blockchain Supply Chain Hackathon  ·  From Farm to Fork — Verified, Transparent, Traceable",
    )

    bullets_slide(prs, "The Problem", [
        "Agricultural supply chains are opaque once produce leaves the farm.",
        ("Food fraud: fake 'organic' claims, mislabelled origin, counterfeits", 1),
        ("No accountability: contamination can't be traced to its source", 1),
        ("Consumer distrust: no reliable way to verify what you buy", 1),
        ("Unfair farmer margins: middlemen obscure true origin and price", 1),
        "Records live in siloed spreadsheets and forgeable paper certificates.",
    ])

    bullets_slide(prs, "Why Existing Systems Fail", [
        "Centralised databases can be edited silently — no tamper-evidence.",
        "Paper certificates are trivially forged or lost.",
        "Siloed data: each actor keeps its own records; no shared truth.",
        "No consumer-facing verification at the point of purchase.",
    ])

    bullets_slide(prs, "Our Solution — AgriChain", [
        "Every batch gets a unique Batch ID (RICE-KONASE-2026-0001).",
        "Each supply-chain event is an immutable, hash-linked block.",
        "Any later edit breaks the chain and is detected instantly.",
        "Compact record + data_hash on-chain; full data off-chain (scalable).",
        "Consumers scan a QR code to verify the full journey.",
        "AI scores risk + flags IoT anomalies; documents are hash-anchored.",
    ])

    image_slide(prs, "System Architecture", ARCHITECTURE_PNG,
                "One authoritative ledger in the FastAPI backend; UI talks only over HTTP.")

    bullets_slide(prs, "Technology Stack", [
        "Backend / API: FastAPI, Uvicorn, Pydantic v2",
        "Blockchain: custom Proof-of-Work engine (SHA-256, Python)",
        "Off-chain store: SQLite (WAL mode)",
        "Frontend: Streamlit + Plotly (9 role-based pages)",
        "AI/ML: scikit-learn — RandomForest (risk) + IsolationForest (anomaly)",
        "IoT: NumPy sensor simulation  ·  QR: qrcode + Pillow",
        "Testing: pytest + httpx (21 tests)  ·  Deck: python-pptx",
    ])

    bullets_slide(prs, "How It Works — Batch Lifecycle", [
        "1. Register batch → HARVEST block mined + off-chain batch row.",
        "2. Record events: QUALITY_CHECK → TRANSPORT → WAREHOUSE → PROCESSING → DISTRIBUTION → RETAIL.",
        "3. Stream IoT telemetry → IsolationForest flags anomalies.",
        "4. Generate QR → consumer scans → verify page.",
        "5. Consumer sees Origin / Quality / Supply-chain / Blockchain ✓.",
        "Each event is its own mined block; batch status advances automatically.",
    ])

    bullets_slide(prs, "Blockchain Design", [
        "hash = SHA-256(index, timestamp, transactions, previous_hash, nonce).",
        ("Stored hash is excluded from its own digest; canonical JSON (sort_keys).", 1),
        "Proof-of-Work: mine until hash has DIFFICULTY (=3) leading zeros.",
        "Previous-hash linking chains every block to its predecessor.",
        "is_chain_valid() re-derives every hash + checks all links.",
        ("Editing a stored tx without re-mining → validation fails.", 1),
    ])

    bullets_slide(prs, "Live Demo Flow", [
        "1. Register a Rice batch on the dashboard.",
        "2. Add QUALITY_CHECK + TRANSPORT events (watch blocks grow).",
        "3. Stream IoT sensors with an injected 89°C spike → anomaly flagged.",
        "4. Show the Blockchain Explorer (hashes, nonces, links).",
        "5. Generate + scan the QR → Consumer verify page shows VERIFIED ✓.",
        "6. Upload a certificate → hash recorded on-chain.",
        "7. Then the money shot →",
    ])

    bullets_slide(prs, "The Money Shot — Tamper Detection", [
        "Silently edit a stored value WITHOUT re-mining:",
        ("quantity_kg:  2500  →  9,999,999", 1),
        "Re-verify the chain:",
        ("chain_valid:  TRUE  →  FALSE", 1),
        "⚠ INTEGRITY COMPROMISED — the recomputed hash no longer matches.",
        "Reset restores a clean, valid ledger. Reproduced offline in tamper_demo.py.",
    ])

    image_slide(prs, "Explainable AI + IoT", FEATURE_IMPORTANCE_PNG,
                "Risk is driven mostly by temperature (0.44), delay (0.31), quality (0.17).")

    bullets_slide(prs, "Results", [
        "36-block seeded chain · 5 batches · 35 events · chain valid ✓.",
        "21/21 automated tests passing (blockchain + API end-to-end).",
        "Tamper detection demonstrably flips chain_valid to false on any edit.",
        "Document verification: MATCH vs MODIFIED on a single changed byte.",
        "All mandatory features + bonus (AI, IoT, documents) delivered.",
    ])

    bullets_slide(prs, "Security — Q&A Highlights", [
        "Tamper prevention: content hashing + PoW + prev-hash links.",
        "On-chain = compact record + data_hash; off-chain = full data (bound by hash).",
        "Documents: SHA-256 anchored on-chain → re-hash to detect edits.",
        "Fake QR: only a pointer; the blockchain is the trust anchor.",
        "Honest limits: single-node ledger + self-asserted identity (roadmap).",
    ])

    bullets_slide(prs, "Future Scope", [
        "Permissioned multi-node network (e.g. Hyperledger Fabric) for consensus.",
        "Per-actor digital signatures → non-repudiable events.",
        "Real IoT hardware (LoRaWAN cold-chain sensors).",
        "Mobile-first consumer PWA + multilingual QR landing pages.",
        "Smart-contract-style automated compliance rules.",
    ])

    title_slide(prs, "Thank You",
                "AgriChain — Farm to Fork, Verified on the Blockchain",
                "Questions?  ·  Live demo + code + docs available")

    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(DECK_PATH)
    return DECK_PATH


def main() -> int:
    path = build()
    print(f"Wrote {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
